"""
文档解析 —— 把各种格式变成「带位置信息的纯文本」
====================================================================
C7 要求覆盖 txt/md/pdf/docx/xlsx/pptx/epub/code/html。

设计要点：

① **返回的不是一坨字符串，是一串带位置的片段。**
   因为搜索结果要能定位到"第几页""第几张幻灯片""哪个工作表"，
   一开始就丢掉位置信息，后面再想加就得重新解析全部文件。

② **每个解析器都可能失败，失败不能拖垮整批。**
   一个损坏的 PDF 不该让 1 万个文件的批处理中断。

③ **重量级依赖延迟导入。** PyMuPDF 之类的库导入要几百毫秒，
   写在模块顶部会让引擎每次启动都慢一截，而且没装这些包时
   整个模块都 import 不进来。
"""

from __future__ import annotations

import csv
import json
import logging
import re
from collections.abc import Iterator
from dataclasses import dataclass
from pathlib import Path

log = logging.getLogger("synorive.parse")

#: 单个文件最多提取多少字符。超大文件（几百 MB 的日志）全读进来会撑爆内存，
#: 而且没人会去搜一个 200MB 日志文件的第 800 万行。
MAX_CHARS = 4_000_000


@dataclass
class TextSegment:
    """一段带位置信息的文本。位置信息决定了搜索结果能不能精确定位。"""

    text: str
    #: 来源通道：body 正文 / title 标题 / ocr 图片文字 / transcript 语音
    channel: str = "body"
    page: int | None = None
    #: 电子表格的工作表名 / PPT 的幻灯片标题 / PDF 的章节
    section: str | None = None


@dataclass
class ParsedDoc:
    segments: list[TextSegment]
    title: str
    #: 页数 / 幻灯片数 / 工作表数
    page_count: int | None = None
    language: str | None = None
    author: str | None = None
    #: 解析过程中的非致命问题，写进 item 的 error 字段给用户看
    warnings: list[str] | None = None

    @property
    def full_text(self) -> str:
        return "\n\n".join(s.text for s in self.segments if s.text.strip())

    @property
    def char_count(self) -> int:
        return sum(len(s.text) for s in self.segments)


class ParseError(RuntimeError):
    pass


# ── 扩展名 → 解析器 ─────────────────────────────────────────

PLAIN_EXT = {
    ".txt", ".md", ".markdown", ".rst", ".log", ".ini", ".cfg", ".conf",
    ".yaml", ".yml", ".toml", ".env", ".properties", ".srt", ".vtt", ".ass",
}

CODE_EXT = {
    ".py", ".js", ".jsx", ".ts", ".tsx", ".java", ".kt", ".c", ".h", ".cpp",
    ".hpp", ".cs", ".go", ".rs", ".rb", ".php", ".swift", ".sh", ".ps1",
    ".sql", ".r", ".m", ".lua", ".vue", ".svelte", ".scss", ".css", ".less",
}

SUPPORTED_EXT = PLAIN_EXT | CODE_EXT | {
    ".pdf", ".docx", ".xlsx", ".xlsm", ".pptx", ".epub",
    ".html", ".htm", ".mhtml", ".json", ".csv", ".tsv",
}


def can_parse(path: Path) -> bool:
    return path.suffix.lower() in SUPPORTED_EXT


def parse(path: Path) -> ParsedDoc:
    """按扩展名分发。任何解析器抛出的异常都包成 ParseError。"""
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            return _parse_pdf(path)
        if ext == ".docx":
            return _parse_docx(path)
        if ext in (".xlsx", ".xlsm"):
            return _parse_xlsx(path)
        if ext == ".pptx":
            return _parse_pptx(path)
        if ext == ".epub":
            return _parse_epub(path)
        if ext in (".html", ".htm", ".mhtml"):
            return _parse_html(path)
        if ext == ".json":
            return _parse_json(path)
        if ext in (".csv", ".tsv"):
            return _parse_csv(path)
        if ext in PLAIN_EXT or ext in CODE_EXT:
            return _parse_plain(path)
    except ParseError:
        raise
    except Exception as e:  # noqa: BLE001
        raise ParseError(f"{path.name} 解析失败：{type(e).__name__}: {e}") from e

    raise ParseError(f"不支持的格式：{ext}")


# ── 纯文本 / 代码 ───────────────────────────────────────────


def _read_text(path: Path) -> str:
    """
    读文本并猜编码。中文文件在国内极大概率是 GBK/GB18030 而不是 UTF-8，
    硬按 UTF-8 读会抛异常或读出乱码 —— 乱码更麻烦，因为它不报错，
    只会让这个文件在搜索里永远命中不了。
    """
    raw = path.read_bytes()[: MAX_CHARS * 4]

    # BOM 优先
    for bom, enc in ((b"\xef\xbb\xbf", "utf-8-sig"), (b"\xff\xfe", "utf-16-le"), (b"\xfe\xff", "utf-16-be")):
        if raw.startswith(bom):
            return raw.decode(enc, errors="replace")

    for enc in ("utf-8", "gb18030", "big5", "shift_jis", "latin-1"):
        try:
            text = raw.decode(enc)
            # UTF-8 解成功基本就是对的；其它编码要检查一下有没有解出大量替换符
            if enc != "utf-8" and text.count("�") > len(text) * 0.01:
                continue
            return text
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def display_title(path: Path) -> str:
    """
    结果列表里显示的标题。

    代码和纯文本文件**带扩展名**：`text.py` 去掉扩展名就剩 `text`，
    库里有 db.py / schema.sql / index.ts 一堆，全显示成 db、schema、index，
    根本分不清是哪个。而 Office 文档的扩展名没有信息量（`年度报告.docx` → `年度报告` 更好看）。
    """
    if path.suffix.lower() in CODE_EXT or path.suffix.lower() in PLAIN_EXT:
        return path.name
    return path.stem


def _parse_plain(path: Path) -> ParsedDoc:
    text = _read_text(path)[:MAX_CHARS]
    return ParsedDoc(segments=[TextSegment(text=text)], title=display_title(path))


def _parse_json(path: Path) -> ParsedDoc:
    text = _read_text(path)[:MAX_CHARS]
    try:
        data = json.loads(text)
        # 重新格式化：紧凑的 JSON 一行几万字符，分块时切不出有意义的边界
        pretty = json.dumps(data, ensure_ascii=False, indent=2)[:MAX_CHARS]
        return ParsedDoc(segments=[TextSegment(text=pretty)], title=display_title(path))
    except json.JSONDecodeError:
        # 不是合法 JSON 就当纯文本，别因为格式错误就整个索引不了
        return ParsedDoc(
            segments=[TextSegment(text=text)],
            title=display_title(path),
            warnings=["JSON 格式不合法，按纯文本索引"],
        )


def _parse_csv(path: Path) -> ParsedDoc:
    text = _read_text(path)
    delim = "\t" if path.suffix.lower() == ".tsv" else ","
    rows = list(csv.reader(text.splitlines()[:50_000], delimiter=delim))
    if not rows:
        return ParsedDoc(segments=[], title=path.stem)

    header = rows[0]
    segs: list[TextSegment] = [TextSegment(text=" | ".join(header), channel="title")]
    # 每行拼成「列名: 值」，这样搜"金额 1200"这种才搜得到
    buf: list[str] = []
    for r in rows[1:]:
        buf.append(" ".join(f"{h}: {v}" for h, v in zip(header, r) if v.strip()))
        if len(buf) >= 200:
            segs.append(TextSegment(text="\n".join(buf)))
            buf = []
    if buf:
        segs.append(TextSegment(text="\n".join(buf)))
    return ParsedDoc(segments=segs, title=path.stem, page_count=len(rows))


# ── PDF ─────────────────────────────────────────────────────
#
# L3 分节：论文按 Abstract/Method/Result 分节索引，搜「实验方法」
# 直接定位到那一节 —— 这是"极速搜索大量文献"里让搜索结果不再是
# "在这篇 20 页论文里翻"，而是"就是这一段"的关键一步。
#
# 判据刻意保守，和这个项目一贯的"宁可漏也不要脏"是同一条原则（图谱实体抽取
# 踩过反面教材：宽松判据把"索引""宋体"这类词误判成人名）。
# 一行文字要同时满足：① 去掉可选的编号前缀和末尾冒号后，
# ② **整行**（不是行的一部分）精确匹配已知的论文章节名，才算标题。
# 只用"这行很短又是黑体"这种版面特征来判会把小标题、图注、作者单位
# 一起认成章节标题；只匹配已知名称虽然会漏掉一些用词生僻的论文，
# 但绝不会把正文一句话误判成新章节，检索结果不会因此"文不对题"。

#: 章节别名 → 归一化后的粗粒度章节名。中英文都收，同一个概念只留一个显示名，
#: 不然搜索结果里 "Method" 和 "Methodology" 被当成两个不同的地方
_SECTION_ALIASES: dict[str, str] = {
    "abstract": "Abstract", "摘要": "Abstract",
    "introduction": "Introduction", "引言": "Introduction", "绪论": "Introduction",
    "related work": "Related Work", "related works": "Related Work", "相关工作": "Related Work",
    "background": "Background", "背景": "Background", "预备知识": "Background",
    "method": "Method", "methods": "Method", "methodology": "Method",
    "materials and methods": "Method", "approach": "Method", "proposed method": "Method",
    "方法": "Method", "研究方法": "Method",
    "experiment": "Experiments", "experiments": "Experiments", "experimental setup": "Experiments",
    "experimental results": "Results", "实验": "Experiments", "实验设置": "Experiments",
    "result": "Results", "results": "Results", "evaluation": "Results", "结果": "Results",
    "实验结果": "Results", "评估": "Results",
    "discussion": "Discussion", "讨论": "Discussion",
    "conclusion": "Conclusion", "conclusions": "Conclusion", "concluding remarks": "Conclusion",
    "future work": "Conclusion", "结论": "Conclusion", "总结": "Conclusion",
    "acknowledgment": "Acknowledgments", "acknowledgement": "Acknowledgments",
    "acknowledgments": "Acknowledgments", "acknowledgements": "Acknowledgments", "致谢": "Acknowledgments",
    "references": "References", "bibliography": "References", "参考文献": "References",
    "appendix": "Appendix", "appendices": "Appendix", "附录": "Appendix",
}

#: 整行必须匹配：可选数字编号（1. / 3.2 / IV.）+ 章节名 + 可选冒号，前后不能有别的字
_SECTION_HEADING_RE = re.compile(
    r"^(?:(?:\d+(?:\.\d+)*\.?)|(?:[IVXLC]+\.))?\s*("
    + "|".join(re.escape(k) for k in sorted(_SECTION_ALIASES, key=len, reverse=True))
    + r")\s*[:：]?\s*$",
    re.IGNORECASE,
)


def _split_pdf_sections(pages: list[tuple[int, str]]) -> list[TextSegment]:
    """把逐页正文按检测到的章节标题切开，每段带上它所属的章节名。"""
    segs: list[TextSegment] = []
    current_section: str | None = None
    buf: list[str] = []
    buf_page: int | None = None

    def flush() -> None:
        text = "\n".join(buf).strip()
        if text:
            segs.append(TextSegment(text=text, page=buf_page, section=current_section))
        buf.clear()

    for page_no, text in pages:
        for line in text.split("\n"):
            m = _SECTION_HEADING_RE.match(line.strip())
            if m:
                flush()
                current_section = _SECTION_ALIASES[m.group(1).lower()]
                buf_page = page_no
                continue
            if buf_page is None:
                buf_page = page_no
            buf.append(line)
        # 换页时也切一刀，保持"位置信息=页码"的既有语义不受影响
        flush()
        buf_page = None
    flush()
    return segs


def _parse_pdf(path: Path) -> ParsedDoc:
    import fitz  # PyMuPDF，延迟导入

    doc = fitz.open(str(path))
    try:
        pages: list[tuple[int, str]] = []
        total = 0
        empty_pages = 0

        for i, page in enumerate(doc):
            text = page.get_text("text").strip()
            if not text:
                empty_pages += 1
                continue
            pages.append((i + 1, text))
            total += len(text)
            if total > MAX_CHARS:
                break

        segs = _split_pdf_sections(pages)

        meta = doc.metadata or {}
        title = (meta.get("title") or "").strip() or path.stem

        warnings: list[str] = []
        # 整本都没文字 = 扫描件，需要 OCR。这不是错误，是要走另一条流水线。
        if empty_pages == doc.page_count and doc.page_count > 0:
            warnings.append("整本没有可提取的文字，是扫描件 —— 需要开 OCR 才能索引内容")
        elif empty_pages > doc.page_count * 0.5:
            warnings.append(f"{empty_pages}/{doc.page_count} 页没有文字层，可能是扫描件混排")

        return ParsedDoc(
            segments=segs,
            title=title,
            page_count=doc.page_count,
            author=(meta.get("author") or "").strip() or None,
            warnings=warnings or None,
        )
    finally:
        doc.close()


# ── Office ──────────────────────────────────────────────────


def _parse_docx(path: Path) -> ParsedDoc:
    import docx

    d = docx.Document(str(path))
    segs: list[TextSegment] = []
    buf: list[str] = []

    for para in d.paragraphs:
        t = para.text.strip()
        if not t:
            continue
        # 标题单独成段并标记，这样标题命中能加权（D4 的 titleBoost）
        if para.style is not None and str(para.style.name or "").lower().startswith("heading"):
            if buf:
                segs.append(TextSegment(text="\n".join(buf)))
                buf = []
            segs.append(TextSegment(text=t, channel="title", section=t))
        else:
            buf.append(t)
    if buf:
        segs.append(TextSegment(text="\n".join(buf)))

    # 表格内容也要索引，很多关键信息就在表格里
    for ti, table in enumerate(d.tables):
        rows = []
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            segs.append(TextSegment(text="\n".join(rows), section=f"表格 {ti + 1}"))

    core = d.core_properties
    return ParsedDoc(
        segments=segs,
        title=(core.title or "").strip() or path.stem,
        author=(core.author or "").strip() or None,
    )


def _parse_xlsx(path: Path) -> ParsedDoc:
    import openpyxl

    # read_only + data_only：不加的话大表会把内存吃光，
    # data_only 让公式返回缓存的计算结果而不是公式文本
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    try:
        segs: list[TextSegment] = []
        total = 0
        for ws in wb.worksheets:
            buf: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if cells:
                    buf.append(" | ".join(cells))
                if len(buf) >= 500:
                    segs.append(TextSegment(text="\n".join(buf), section=ws.title))
                    total += sum(len(x) for x in buf)
                    buf = []
                if total > MAX_CHARS:
                    break
            if buf:
                segs.append(TextSegment(text="\n".join(buf), section=ws.title))
                total += sum(len(x) for x in buf)
            if total > MAX_CHARS:
                break
        return ParsedDoc(segments=segs, title=path.stem, page_count=len(wb.worksheets))
    finally:
        wb.close()


def _parse_pptx(path: Path) -> ParsedDoc:
    from pptx import Presentation

    prs = Presentation(str(path))
    segs: list[TextSegment] = []

    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        slide_title = None
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            t = shape.text_frame.text.strip()
            if not t:
                continue
            if slide_title is None:
                slide_title = t.splitlines()[0][:60]
            parts.append(t)
        # 演讲者备注也索引 —— 那里往往才是真正的内容
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                parts.append(f"[备注] {note}")
        if parts:
            segs.append(TextSegment(text="\n".join(parts), page=i, section=slide_title))

    return ParsedDoc(segments=segs, title=path.stem, page_count=len(prs.slides))


def _parse_epub(path: Path) -> ParsedDoc:
    try:
        import ebooklib
        from ebooklib import epub
    except ImportError as e:
        raise ParseError("EPUB 需要 ebooklib，跑一次依赖医生装上") from e

    from html import unescape
    import re

    book = epub.read_epub(str(path))
    segs: list[TextSegment] = []
    total = 0
    for i, item in enumerate(book.get_items_of_type(ebooklib.ITEM_DOCUMENT)):
        html = item.get_content().decode("utf-8", errors="replace")
        text = unescape(re.sub(r"<[^>]+>", " ", html))
        text = re.sub(r"\s+", " ", text).strip()
        if text:
            segs.append(TextSegment(text=text, page=i + 1))
            total += len(text)
        if total > MAX_CHARS:
            break

    titles = book.get_metadata("DC", "title")
    return ParsedDoc(
        segments=segs,
        title=(titles[0][0] if titles else path.stem),
        page_count=len(segs),
    )


# ── 网页 ────────────────────────────────────────────────────


def _parse_html(path: Path) -> ParsedDoc:
    html = _read_text(path)
    return parse_html_string(html, fallback_title=path.stem)


def parse_html_string(html: str, *, fallback_title: str = "", url: str | None = None) -> ParsedDoc:
    """
    从 HTML 里抽正文。C11 存档和本地 .html 文件共用这一条路径。

    trafilatura 会去掉导航、广告、页脚这些噪声 —— 不去的话
    每个网页都会因为共同的导航文字而互相"相似"，语义检索直接废掉。
    """
    import re
    from html import unescape

    title = fallback_title
    m = re.search(r"<title[^>]*>(.*?)</title>", html, re.S | re.I)
    if m:
        title = unescape(re.sub(r"\s+", " ", m.group(1))).strip() or fallback_title

    try:
        import trafilatura

        text = trafilatura.extract(
            html,
            url=url,
            include_comments=False,
            include_tables=True,
            favor_precision=True,
        )
        if text and text.strip():
            return ParsedDoc(
                segments=[TextSegment(text=text.strip()[:MAX_CHARS])],
                title=title,
            )
    except Exception as e:  # noqa: BLE001
        log.debug("trafilatura 抽取失败，回退到粗暴去标签：%s", e)

    # 兜底：粗暴去标签。质量差但总比索引不了强。
    body = re.sub(r"<(script|style|nav|header|footer)[^>]*>.*?</\1>", " ", html, flags=re.S | re.I)
    text = unescape(re.sub(r"<[^>]+>", " ", body))
    text = re.sub(r"\s+", " ", text).strip()
    return ParsedDoc(
        segments=[TextSegment(text=text[:MAX_CHARS])],
        title=title,
        warnings=["正文提取器没抽出内容，用了粗暴去标签，可能含导航噪声"],
    )


#: 目录扫描（手动投喂 + watcher.py 目录监控）共用的垃圾目录名单。
#: 两处都要跳过这些——手动投喂已经在用，监控目录如果不跳过，
#: 一个带 node_modules 的项目文件夹会被 watchdog 每个文件都盯上，
#: 白白注册几万个监听、库里也会混进一堆没意义的第三方代码。
SKIP_DIR_NAMES = frozenset({
    "node_modules", ".git", ".venv", "venv", "__pycache__", ".idea", ".vscode",
    "dist", "build", "out", ".gradle", ".next", "target", ".cache",
})


def iter_supported(root: Path, recursive: bool = True) -> Iterator[Path]:
    """
    遍历目录里所有能处理的文件（文档 + 图片 + 视频 + 音频）。
    跳过隐藏目录和常见的垃圾目录。

    ⚠️ 加新模态时**这里也要加扩展名**。加了图片忘了加视频，
    症状是"拖一个视频文件夹进去，一个都没被索引"，而且不报错。
    """
    from ..analyze.image import SUPPORTED_IMAGE_EXT
    from ..analyze.video import SUPPORTED_AUDIO_EXT, SUPPORTED_VIDEO_EXT
    from .sensitive import sensitive_reason

    media_ext = SUPPORTED_IMAGE_EXT | SUPPORTED_VIDEO_EXT | SUPPORTED_AUDIO_EXT

    skip_dirs = SKIP_DIR_NAMES
    stack = [root]
    while stack:
        d = stack.pop()
        try:
            entries = list(d.iterdir())
        except (PermissionError, OSError):
            continue
        for p in entries:
            if p.is_dir():
                if recursive and p.name not in skip_dirs and not p.name.startswith("."):
                    stack.append(p)
            elif can_parse(p) or p.suffix.lower() in media_ext:
                yield p
            # 🔴 .env / id_rsa 这类密钥文件本来就没有扩展名，`can_parse()`
            # 天然判它"不支持"，会像 .exe/.zip 一样悄悄从遍历结果里消失——
            # 消失了就不会被下面 ingest_paths() 里的敏感文件检查看到，
            # 用户也就永远不知道这个目录里有个 .env。这里额外放一道口子：
            # 即使格式不支持，只要命中敏感文件规则也要让它被看见（然后
            # 在 ingest_paths() 里被明确标成"跳过：敏感文件"，而不是凭空消失）。
            elif sensitive_reason(p) is not None:
                yield p
